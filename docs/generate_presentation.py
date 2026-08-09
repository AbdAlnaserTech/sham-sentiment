"""
توليد عرض PowerPoint — الجزء النظري (مطلوب الدكتور محمد أسامة).

التشغيل: python docs/generate_presentation.py

ملاحظة: أهم شي المخططات — أضفها يدوياً في الشرائح المحددة.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# جذر المشروع ومسار ملف العرض الناتج
ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "عرض_مشروع_تحليل_المشاعر.pptx"

TITLE_AR = "تحليل مشاعر التعليقات متعددة اللغات"
SUBTITLE = "مشروع تخرج — الهندسة المعلوماتية — جامعة الشام"
SUPERVISOR = "إشراف: د. محمد أسامة"
TEAM = "عبد الناصر الحسون · عدي نجار · حسن البكور · محمود السيد · عدنان ناقوح"
YEAR = "2026"

# ألوان العرض
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
BLUE = RGBColor(0x25, 0x63, 0xEB)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _set_rtl(paragraph) -> None:
    """محاذاة الفقرة من اليمين لليسار."""
    paragraph.alignment = PP_ALIGN.RIGHT


def _add_title_slide(prs: Presentation) -> None:
    """إضافة شريحة العنوان الرئيسية."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True

    lines = [TITLE_AR, SUBTITLE, TEAM, SUPERVISOR, YEAR]
    sizes = [36, 16, 13, 14, 14]
    bolds = [True, False, False, False, False]
    colors = [NAVY, GRAY, GRAY, GRAY, GRAY]

    for i, (text, size, bold, color) in enumerate(zip(lines, sizes, bolds, colors)):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        _set_rtl(p)


def _slide(prs: Presentation, title: str, bullets: list[str], *, note: str = "") -> None:
    """إضافة شريحة محتوى بعنوان ونقاط."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # شريط العنوان العلوي
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.7))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    _set_rtl(tp)

    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
        p.space_after = Pt(10)
        _set_rtl(p)

    if note:
        nb = slide.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(8.8), Inches(0.8))
        np = nb.text_frame.paragraphs[0]
        np.text = note
        np.font.size = Pt(14)
        np.font.color.rgb = BLUE
        np.font.italic = True
        _set_rtl(np)


def _diagram_slide(prs: Presentation, title: str, diagram_hint: str) -> None:
    """إضافة شريحة placeholder للمخططات (يُضاف المخطط يدوياً)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.7))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    _set_rtl(tp)

    frame = slide.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.2))
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
    frame.line.color.rgb = BLUE

    hint = slide.shapes.add_textbox(Inches(1.2), Inches(2.8), Inches(7.6), Inches(2))
    hp = hint.text_frame.paragraphs[0]
    hp.text = diagram_hint
    hp.font.size = Pt(20)
    hp.font.color.rgb = GRAY
    hp.alignment = PP_ALIGN.CENTER


def build() -> Path:
    """بناء العرض الكامل وحفظه في docs/."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs)

    _slide(
        prs,
        "فهرس العرض",
        [
            "فكرة عامة عن المشروع",
            "مشكلة البحث",
            "هدف المشروع",
            "المنهجية (Agile / SDLC)",
            "الدراسات السابقة",
            "الجدول الزمني",
            "المتطلبات الوظيفية وغير الوظيفية",
            "دراسة الجدوى الاقتصادية",
            "نمذجة النظام (مخططات)",
            "خوارزميات الذكاء الصنعي",
        ],
    )

    _slide(
        prs,
        "1 — فكرة عامة عن المشروع",
        [
            "منصة ويب لتحليل مشاعر التعليقات (إيجابي · محايد · سلبي).",
            "دعم العربية الفصحى، اللهجة الشامية، والإنجليزية.",
            "إدخال تعليق واحد، دفعة CSV، أو جلب من YouTube / Google Play.",
            "عرض النتائج، الثقة، التحليلات، التنبيهات، وتصدير PDF/Excel.",
        ],
    )

    _slide(
        prs,
        "2 — مشكلة البحث",
        [
            "حجم التعليقات كبير — التحليل اليدوي بطيء وغير عملي.",
            "أغلب الأدوات لا تدعم اللهجة الشامية بشكل جيد.",
            "صعوبة دمج التحليل + التخزين + التقارير في حل واحد.",
            "حاجة المؤسسات لفهم آراء العملاء بسرعة ودقة.",
        ],
    )

    _slide(
        prs,
        "3 — هدف المشروع",
        [
            "بناء نظام تصنيف مشاعر ثلاثي الفئات.",
            "كشف اللغة تلقائياً (en / ar_fusha / ar_shami).",
            "توفير واجهة عربية + صلاحيات مستخدمين.",
            "مقارنة BERT مع TF-IDF وتقييم على بيانات تحقق.",
        ],
    )

    _slide(
        prs,
        "4 — المنهجية المستخدمة",
        [
            "Agile (تطوير تكراري) ضمن SDLC.",
            "التخطيط → المتطلبات → التحليل → التصميم → التنفيذ → الاختبار → التوثيق.",
            "مراجعات دورية مع المشرف.",
            "43 unit test (pytest) قبل التسليم.",
        ],
    )

    _slide(
        prs,
        "5 — الدراسات السابقة",
        [
            "Pang & Lee — أسس Sentiment Analysis.",
            "ASTD — dataset tweets عربية.",
            "BERT / XLM-RoBERTa — Transformers متعددة اللغات.",
            "CAMeLBERT — نماذج عربية.",
            "تمييز مشروعنا: منصة متكاملة + دعm شامي + واجهة عربية.",
        ],
    )

    _diagram_slide(
        prs,
        "6 — الجدول الزمني / Gantt",
        "[ أضف مخطط Gantt هنا ]\n\nتخطيط(3) · متطلبات(4) · تحليل/تصميم(6)\nنماذج(6) · واجهة(5) · اختبار(4)",
    )

    _slide(
        prs,
        "7 — المتطلبات الوظيفية",
        [
            "FR1: تسجيل دخول (admin / analyst / viewer).",
            "FR2: تحليل تعليق واحد + wordcloud.",
            "FR3: تحليل دفعة CSV.",
            "FR4: جلب تعليقات YouTube / Google Play.",
            "FR5: لوحة تحكم + KPIs + تنبيهات.",
            "FR6: تصدير PDF / Excel.",
        ],
    )

    _slide(
        prs,
        "8 — المتطلبات غير الوظيفية",
        [
            "NFR1: أمان — PBKDF2-SHA256.",
            "NFR2: أداء — cache لنموذج BERT.",
            "NFR3: RTL — واجهة عربية.",
            "NFR4: حد أقصى 5000 حرف للتعليق.",
        ],
    )

    _slide(
        prs,
        "9 — دراسة الجدوى الاقتصادية",
        [
            "تكلفة منخفضة — Python مفتوح المصدر.",
            "Streamlit Cloud مجاني للعرض.",
            "لا تراخيص مدفوعة للنماذج الأساسية.",
            "قابل للتوسع commerciaً لاحقاً.",
        ],
    )

    _diagram_slide(
        prs,
        "10 — مخطط السياق (Context Diagram)",
        "[ أضف Context Diagram ]\n\nالمستخدم ↔ المنصة ↔ YouTube/Play\n                    ↔ SQLite ↔ BERT Model",
    )

    _diagram_slide(
        prs,
        "11 — مخطط حالات الاستخدام (Use Case)",
        "[ أضف Use Case Diagram ]\n\nActors: Admin, Analyst, Viewer\nUC: Login, Analyze, Batch, Import, Dashboard",
    )

    _slide(
        prs,
        "12 — وصف حالات الاستخدام",
        [
            "UC1 تسجيل الدخول: يُدخل username/password → يتحقق النظام → dashboard.",
            "UC2 تحليل تعليق: إدخال نص → كشف لغة → BERT → عرض sentiment + confidence.",
            "UC3 دفعة: رفع CSV → loop تحليل → جدول + pie chart.",
            "UC4 جلب حي: URL → fetcher → preview → تحليل.",
        ],
    )

    _diagram_slide(
        prs,
        "13 — مخطط النشاط (Activity)",
        "[ أضف Activity Diagram ]\n\nتعليق → preprocessing → BERT → نتيجة → (حفظ؟) → export",
    )

    _diagram_slide(
        prs,
        "14 — مخطط التتابع (Sequence)",
        "[ أضف Sequence Diagram ]\n\nUser → Streamlit → Predictor → BERT → SQLite → Response",
    )

    _slide(
        prs,
        "15 — خوارزميات الذكاء الصنعي (1)",
        [
            "BERT / XLM-RoBERTa: Transformer مُدرَّب مسبقاً، fine-tune للمشاعر.",
            "Input: tokens → Output: probabilities (neg/neu/pos).",
            "Ensemble: XLM-R + CAMeLBERT مع fallbacks.",
            "Default في الواجهة — أعلى دقة (65.3%).",
        ],
    )

    _slide(
        prs,
        "16 — شرح الخوارزميات (2)",
        [
            "TF-IDF + Logistic Regression: baseline سريع.",
            "GridSearchCV + CalibratedClassifierCV للثقة.",
            "language.py: كشف en / ar_fusha / ar_shami (قواعد + hints).",
            "preprocessing.py: normalize_arabic, stopwords, emoji.",
        ],
    )

    _slide(
        prs,
        "17 — النتائج (مختصر)",
        [
            "validation_comments.csv — 513 تعليق.",
            "BERT: Accuracy 65.3% | Macro-F1 0.63",
            "TF-IDF: Accuracy 53.8% | Macro-F1 0.53",
            "أفضل في positive/negative — neutral يحتاج تحسين.",
        ],
        note="[ أضف Confusion Matrix كشكل ]",
    )

    _slide(
        prs,
        "شكراً لكم",
        [
            "أسئلة؟",
            "",
            SUPERVISOR,
            TEAM,
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    print(f"Saved: {build()}")
